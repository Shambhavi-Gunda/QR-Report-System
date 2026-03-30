import pytesseract
import cv2
import numpy as np
import os
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io
from PyPDF2 import PdfReader

# SET TESSERACT PATH
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"


# --------------------------------------------------
# IMAGE PREPROCESSING (IMPROVED)
# --------------------------------------------------
def preprocess_image_from_array(img):
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

    # Upscale image (critical for OCR accuracy)
    gray = cv2.resize(gray, None, fx=2, fy=2, interpolation=cv2.INTER_CUBIC)

    # Noise removal
    gray = cv2.bilateralFilter(gray, 9, 75, 75)

    # Adaptive threshold (better than fixed threshold)
    thresh = cv2.adaptiveThreshold(
        gray, 255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        31, 2
    )

    return thresh


# --------------------------------------------------
# OCR CORE FUNCTION
# --------------------------------------------------
def ocr_from_array(img):
    try:
        processed = preprocess_image_from_array(img)

        text = pytesseract.image_to_string(
            processed,
            config='--oem 3 --psm 6 -l eng'
        )

        return text.strip()

    except Exception as e:
        print("OCR error:", e)
        return ""


# --------------------------------------------------
# IMAGE OCR
# --------------------------------------------------
def extract_text_from_image(image_path):
    try:
        img = cv2.imread(image_path)

        if img is None:
            print("Failed to read image:", image_path)
            return ""

        text = ocr_from_array(img)
        return text

    except Exception as e:
        print("Image OCR error:", e)
        return ""


# --------------------------------------------------
# CHECK IF PDF IS SCANNED
# --------------------------------------------------
def is_scanned_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)

        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                return False  # text-based PDF

        return True  # scanned PDF

    except Exception as e:
        print("PDF check error:", e)
        return True


# --------------------------------------------------
# OCR PDF (ONLY WHEN NEEDED)
# --------------------------------------------------
def ocr_pdf(pdf_path):
    text = ""

    try:
        images = convert_from_path(
            pdf_path,
            dpi=300,
            poppler_path=r"C:\poppler\poppler-25.12.0\Library\bin"
        )

        for img in images:
            img_np = np.array(img)
            page_text = ocr_from_array(img_np)
            text += page_text + "\n"

        return text.strip()

    except Exception as e:
        print("PDF OCR error:", e)
        return ""


# --------------------------------------------------
# SMART PDF EXTRACTION
# --------------------------------------------------
def extract_text_from_pdf(pdf_path):
    text = ""

    try:
        reader = PdfReader(pdf_path)

        # Step 1: Try direct extraction
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted

        # Step 2: If no text → use OCR
        if not text.strip():
            print("Using OCR fallback for scanned PDF")
            return ocr_pdf(pdf_path)

        return text.strip()

    except Exception as e:
        print("PDF extraction error:", e)
        return ""


# --------------------------------------------------
# DOCX TEXT
# --------------------------------------------------
def extract_text_from_docx(docx_path):
    try:
        doc = Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()

    except Exception as e:
        print("DOCX error:", e)
        return ""


# --------------------------------------------------
# PPTX TEXT + IMAGE OCR
# --------------------------------------------------
def extract_text_from_pptx(pptx_path):
    text = ""

    try:
        prs = Presentation(pptx_path)

        for slide in prs.slides:
            for shape in slide.shapes:

                # Extract text
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

                # OCR images inside PPT
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob

                    img = Image.open(io.BytesIO(image_bytes))
                    img_np = np.array(img)

                    ocr_text = ocr_from_array(img_np)
                    text += ocr_text + "\n"

        return text.strip()

    except Exception as e:
        print("PPTX error:", e)
        return ""


# --------------------------------------------------
# TXT FILE
# --------------------------------------------------
def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, "r", encoding="utf-8") as f:
            text = f.read()

        return text.strip()

    except Exception as e:
        print("TXT error:", e)
        return ""


# --------------------------------------------------
# MAIN ROUTER
# --------------------------------------------------
def extract_text(file_path):
    ext = file_path.split('.')[-1].lower()

    print(f"Processing file: {file_path} | Type: {ext}")

    if ext in ['png', 'jpg', 'jpeg']:
        return extract_text_from_image(file_path)

    elif ext == 'pdf':
        return extract_text_from_pdf(file_path)

    elif ext == 'docx':
        return extract_text_from_docx(file_path)

    elif ext == 'pptx':
        return extract_text_from_pptx(file_path)

    elif ext == 'txt':
        return extract_text_from_txt(file_path)

    else:
        print("Unsupported file type:", ext)
        return ""